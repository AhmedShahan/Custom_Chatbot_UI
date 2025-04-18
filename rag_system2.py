import numpy as np
import faiss
import json
import uuid
import torch
import pandas as pd
import re
from transformers import RobertaModel, RobertaTokenizer
from typing import Dict, List, Tuple, Optional, Union, Any
from dataclasses import dataclass, asdict, field
from tenacity import retry, stop_after_attempt, wait_exponential
from unstructured.partition.pdf import partition_pdf
from unstructured.partition.pptx import partition_pptx
from unstructured.partition.docx import partition_docx
from unstructured.partition.doc import partition_doc
from unstructured.documents.elements import Title, NarrativeText, Table, Element
from langchain.llms import Ollama
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from collections import Counter
import spacy
import nltk
from nltk.tokenize import sent_tokenize
from string import Template
import os
import subprocess
import platform
import shutil
from pptx import Presentation
import tempfile
from pathlib import Path
import time
import io
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from PyPDF2 import PdfWriter, PdfReader
from fpdf import FPDF
from reportlab.lib import colors
from reportlab.lib.pagesizes import landscape
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
from reportlab.lib.units import inch

# Download necessary NLTK resources
nltk.download('punkt', quiet=True)
nltk.download('stopwords', quiet=True)
nltk.download('wordnet', quiet=True)

# Load spaCy model for NER and keyword extraction
try:
    nlp = spacy.load("en_core_web_sm")
except:
    import os
    os.system("python -m spacy download en_core_web_sm")
    nlp = spacy.load("en_core_web_sm")

@dataclass
class DocumentEntry:
    id: str
    title: str
    text: str
    title_embedding: np.ndarray
    text_embedding: np.ndarray
    keywords: List[str] = field(default_factory=list)
    entities: Dict[str, List[str]] = field(default_factory=dict)
    metadata: Dict = field(default_factory=dict)

class DocumentEntry:
    id: str
    title: str
    text: str
    title_embedding: np.ndarray
    text_embedding: np.ndarray
    keywords: List[str] = field(default_factory=list)
    entities: Dict[str, List[str]] = field(default_factory=dict)
    metadata: Dict = field(default_factory=dict)

class ParallelVectorStore:
    def __init__(self, embedding_dim: int = 768):
        self.title_index = faiss.IndexFlatL2(embedding_dim)
        self.text_index = faiss.IndexFlatL2(embedding_dim)
        self.documents: Dict[str, DocumentEntry] = {}
        self.tfidf_vectorizer = None
        self.tfidf_matrix = None
        self.document_ids = []

    def add_document(self, title: str, text: str, 
                    title_embedding: np.ndarray, text_embedding: np.ndarray,
                    keywords: List[str] = None, entities: Dict[str, List[str]] = None,
                    metadata: Dict = None) -> str:
        doc_id = str(uuid.uuid4())
        title_embedding = np.array(title_embedding, dtype=np.float32).reshape(1, -1)
        text_embedding = np.array(text_embedding, dtype=np.float32).reshape(1, -1)

        self.title_index.add(title_embedding)
        self.text_index.add(text_embedding)

        self.documents[doc_id] = DocumentEntry(
            id=doc_id,
            title=title,
            text=text,
            title_embedding=title_embedding.flatten(),
            text_embedding=text_embedding.flatten(),
            keywords=keywords or [],
            entities=entities or {},
            metadata=metadata or {}
        )
        
        self.document_ids.append(doc_id)
        self._rebuild_tfidf()
        return doc_id

    def _rebuild_tfidf(self):
        """Rebuild TF-IDF matrix with all documents"""
        corpus = [f"{doc.title} {doc.text}" for doc in self.documents.values()]
        
        if not corpus:
            print("No documents in corpus, skipping TF-IDF rebuild")
            return
            
        doc_count = len(corpus)
        print(f"Building TF-IDF matrix with {doc_count} documents")
        
        try:
            if doc_count == 1:
                self.tfidf_vectorizer = TfidfVectorizer(
                    max_df=1.0,
                    min_df=1,
                    stop_words='english'
                )
            else:
                min_df_value = 1 if doc_count < 3 else 2
                self.tfidf_vectorizer = TfidfVectorizer(
                    max_df=0.95, 
                    min_df=min_df_value, 
                    stop_words='english'
                )
                
            self.tfidf_matrix = self.tfidf_vectorizer.fit_transform(corpus)
            print(f"TF-IDF matrix built successfully with shape {self.tfidf_matrix.shape}")
        except Exception as e:
            print(f"Error building TF-IDF matrix: {str(e)}")
            self.tfidf_vectorizer = None
            self.tfidf_matrix = None

    def search(self, query_embedding: np.ndarray, k: int = 5) -> List[DocumentEntry]:
        query_embedding = np.array(query_embedding, dtype=np.float32).reshape(1, -1)
        _, title_indices = self.title_index.search(query_embedding, k)
        _, text_indices = self.text_index.search(query_embedding, k)
        combined_indices = set(title_indices[0].tolist() + text_indices[0].tolist())
        return [self.documents[list(self.documents.keys())[idx]] for idx in combined_indices]
    
    def tfidf_search(self, query: str, k: int = 5) -> List[DocumentEntry]:
        if self.tfidf_vectorizer is None or self.tfidf_matrix is None:
            print("TF-IDF search not available, returning empty results")
            return []
            
        try:
            query_vec = self.tfidf_vectorizer.transform([query])
            similarity_scores = cosine_similarity(query_vec, self.tfidf_matrix).flatten()
            top_indices = similarity_scores.argsort()[-k:][::-1]
            return [self.documents[self.document_ids[idx]] for idx in top_indices]
        except Exception as e:
            print(f"Error in TF-IDF search: {str(e)}")
            return []

    def keyword_search(self, keywords: List[str], k: int = 5) -> List[DocumentEntry]:
        keyword_scores = {}
        
        for doc_id, doc in self.documents.items():
            score = 0
            doc_keywords = set(doc.keywords)
            doc_text = f"{doc.title.lower()} {doc.text.lower()}"
            
            for keyword in keywords:
                if keyword.lower() in doc_keywords:
                    score += 3
                score += doc_text.count(keyword.lower())
            
            keyword_scores[doc_id] = score
        
        top_docs = sorted(keyword_scores.items(), key=lambda x: x[1], reverse=True)[:k]
        return [self.documents[doc_id] for doc_id, _ in top_docs]

    def save(self, path_prefix: str):
        faiss.write_index(self.title_index, f"{path_prefix}_title_vectors.faiss")
        faiss.write_index(self.text_index, f"{path_prefix}_text_vectors.faiss")
        documents_to_save = {}
        for doc_id, doc_entry in self.documents.items():
            doc_dict = asdict(doc_entry)
            doc_dict['title_embedding'] = doc_dict['title_embedding'].tolist()
            doc_dict['text_embedding'] = doc_dict['text_embedding'].tolist()
            documents_to_save[doc_id] = doc_dict
        with open(f"{path_prefix}_documents.json", 'w') as f:
            json.dump(documents_to_save, f)

    def print_all_documents(self):
        """Print all chunked information in the vector store."""
        if not self.documents:
            print("No chunked information found in the vector store.")
            return
        
        print("\n=== Chunked Information in Vector Store ===")
        for doc_id, doc in self.documents.items():
            print(f"\nChunk ID: {doc_id}")
            print(f"Text (Full Content, {len(doc.text)} chars):\n{doc.text}")
            print(f"Keywords: {', '.join(doc.keywords)}")
            print(f"Entities: {doc.entities}")
            print(f"Metadata: {doc.metadata}")
            print("=" * 80)

    def print_chunk_by_id(self, doc_id: str):
        """Print chunked information for a specific document ID."""
        if doc_id not in self.documents:
            print(f"No chunk found with ID: {doc_id}")
            return
        doc = self.documents[doc_id]
        print(f"\n=== Chunked Information for ID: {doc_id} ===")
        print(f"Chunk ID: {doc_id}")
        print(f"Text (Full Content, {len(doc.text)} chars):\n{doc.text}")
        print(f"Keywords: {', '.join(doc.keywords)}")
        print(f"Entities: {doc.entities}")
        print(f"Metadata: {doc.metadata}")
        print("=" * 80)

    @classmethod
    def load(cls, path_prefix: str):
        store = cls()
        store.title_index = faiss.read_index(f"{path_prefix}_title_vectors.faiss")
        store.text_index = faiss.read_index(f"{path_prefix}_text_vectors.faiss")
        with open(f"{path_prefix}_documents.json", 'r') as f:
            documents_dict = json.load(f)
        store.documents = {}
        store.document_ids = []
        for k, v in documents_dict.items():
            v['title_embedding'] = np.array(v['title_embedding'], dtype=np.float32)
            v['text_embedding'] = np.array(v['text_embedding'], dtype=np.float32)
            store.documents[k] = DocumentEntry(**v)
            store.document_ids.append(k)
        store._rebuild_tfidf()
        return store

class ResponseTemplate:
    def __init__(self):
        self.templates = {
            "definition": Template("$entity refers to $definition"),
            "summary": Template("Here's a summary of $topic: $summary"),
            "list": Template("Here are the key points about $topic:\n$items"),
            "comparison": Template("Comparing $entity1 and $entity2:\n$comparison"),
            "fallback": Template("Based on the provided information: $content")
        }
        
    def add_template(self, template_type: str, template_string: str):
        self.templates[template_type] = Template(template_string)
        
    def generate_response(self, template_type: str, **kwargs) -> str:
        if template_type not in self.templates:
            template_type = "fallback"
        return self.templates[template_type].safe_substitute(**kwargs)

class RuleBasedExtractor:
    def __init__(self):
        self.entity_patterns = {}
        self.relation_patterns = {}
        
    def add_entity_pattern(self, entity_type: str, pattern: str):
        self.entity_patterns[entity_type] = re.compile(pattern, re.IGNORECASE)
        
    def add_relation_pattern(self, relation_type: str, pattern: str):
        self.relation_patterns[relation_type] = re.compile(pattern, re.IGNORECASE)
        
    def extract_entities(self, text: str) -> Dict[str, List[str]]:
        results = {}
        for entity_type, pattern in self.entity_patterns.items():
            matches = pattern.findall(text)
            if matches:
                results[entity_type] = matches
        doc = nlp(text)
        for ent in doc.ents:
            entity_type = ent.label_
            if entity_type not in results:
                results[entity_type] = []
            results[entity_type].append(ent.text)
        return results
        
    def extract_relations(self, text: str) -> Dict[str, List[str]]:
        results = {}
        for relation_type, pattern in self.relation_patterns.items():
            matches = pattern.findall(text)
            if matches:
                results[relation_type] = matches
        return results

class ExtractiveAnswerGenerator:
    @staticmethod
    def extract_sentences(text: str) -> List[str]:
        if not text or not isinstance(text, str):
            return []
        try:
            sentences = sent_tokenize(text)
            return [s for s in sentences if s.strip()]
        except Exception as e:
            print(f"Error extracting sentences: {str(e)}")
            return []
    
    @staticmethod
    def rank_sentences(query: str, sentences: List[str]) -> List[Tuple[str, float]]:
        if not sentences or not query:
            return []
        sentences = [s for s in sentences if s and s.strip()]
        if not sentences:
            return []
        try:
            vectorizer = TfidfVectorizer(stop_words='english')
            sentence_vectors = vectorizer.fit_transform(sentences)
            query_vector = vectorizer.transform([query])
            similarities = cosine_similarity(query_vector, sentence_vectors).flatten()
            return [(sentence, score) for sentence, score in zip(sentences, similarities)]
        except Exception as e:
            print(f"Error ranking sentences: {str(e)}")
            return [(sentence, 0.0) for sentence in sentences]
    
    def generate_answer(self, query: str, documents: List[DocumentEntry], 
                        max_sentences: int = 3) -> str:
        if not documents or not query:
            return "No relevant information found."
        all_sentences = []
        for doc in documents:
            if not hasattr(doc, 'text') or not doc.text:
                continue
            try:
                sentences = self.extract_sentences(doc.text)
                all_sentences.extend(sentences)
            except Exception as e:
                print(f"Error processing document for sentences: {str(e)}")
                continue
        if not all_sentences:
            return "No relevant information found in the documents."
        try:
            ranked_sentences = self.rank_sentences(query, all_sentences)
            if not ranked_sentences:
                return "Couldn't extract relevant sentences from the documents."
            top_sentences = sorted(ranked_sentences, key=lambda x: x[1], reverse=True)
            top_sentences = top_sentences[:min(max_sentences, len(top_sentences))]
            if not top_sentences:
                return "No relevant information found."
            answer = " ".join([sentence for sentence, _ in top_sentences])
            return answer
        except Exception as e:
            print(f"Error generating extractive answer: {str(e)}")
            return "Encountered an issue while generating an answer from the document."

class RAGSystem:
    def __init__(self, model_name: str = "deepseek-r1:14b", use_llm: bool = True, method: str = "hybrid"):
        self.vector_store = ParallelVectorStore()
        self.templates = ResponseTemplate()
        self.extractor = RuleBasedExtractor()
        self.extractive_generator = ExtractiveAnswerGenerator()
        self.use_llm = use_llm
        self.method = method
        self.model_name = model_name
        self.chunk_count = 0  # Initialize chunk counter
        
        self.developer_info = {
            "name": "Shahan Ahmed",
            "title": "Data Scientist",
            "company": "Startsmartz Technologies LLC",
            "github": "https://github.com/AhmedShahan",
            "researchgate": "https://www.researchgate.net/profile/Shahan-Ahmed-2?ev=hdr_xprf",
            "email": "shahan.ahmed001@gmail.com"
        }

        self.device = torch.device('cuda' if torch.cuda.is_available() else 'cpu')
        self.tokenizer = RobertaTokenizer.from_pretrained('roberta-base')
        self.embedding_model = RobertaModel.from_pretrained('roberta-base').to(self.device)
        self.embedding_model.eval()

        if use_llm:
            try:
                self.generation_model = Ollama(model=model_name)
                print(f"Initialized LLM using Ollama with model: {model_name}")
            except Exception as e:
                print(f"Warning: Failed to initialize Ollama LLM: {str(e)}")
                self.generation_model = None
                self.use_llm = False
                print("Falling back to non-LLM methods only")
            
        self._initialize_extraction_patterns()
        
    def _initialize_extraction_patterns(self):
        self.extractor.add_entity_pattern("email", r'[\w\.-]+@[\w\.-]+\.\w+')
        self.extractor.add_entity_pattern("phone", r'\+?[\d\s-]{10,}')
        self.extractor.add_entity_pattern("url", r'https?://[\w\.-/]+')
        self.extractor.add_entity_pattern("date", r'\d{1,2}[/-]\d{1,2}[/-]\d{2,4}')
        self.extractor.add_entity_pattern("definition", r'(\w+)\s+is\s+defined\s+as\s+([^.]+)')
        self.extractor.add_entity_pattern("comparison", r'compared\s+to\s+(\w+)[,\s]+(\w+)\s+([^.]+)')

    def get_embedding(self, text: str) -> np.ndarray:
        inputs = self.tokenizer(
            text,
            max_length=512,
            padding=True,
            truncation=True,
            return_tensors="pt"
        ).to(self.device)
        with torch.no_grad():
            outputs = self.embedding_model(**inputs)
            embeddings = outputs.last_hidden_state[:, 0, :].cpu().numpy()
        return embeddings[0]

    def extract_keywords(self, text: str, max_keywords: int = 10) -> List[str]:
        doc = nlp(text)
        keywords = []
        keywords.extend([ent.text.lower() for ent in doc.ents])
        for chunk in doc.noun_chunks:
            keywords.append(chunk.text.lower())
        keywords = [k for k in keywords if len(k) > 2]
        counter = Counter(keywords)
        return [word for word, _ in counter.most_common(max_keywords)]

    def _convert_table_to_text(self, table: Table) -> str:
        try:
            if hasattr(table, 'metadata') and table.metadata and 'text_as_html' in table.metadata:
                tables = pd.read_html(table.metadata['text_as_html'])
                if not tables:
                    return str(table)
                df = tables[0]
            else:
                data = table.metadata.get('data', []) if hasattr(table, 'metadata') and table.metadata else []
                if not data:
                    return str(table)
                df = pd.DataFrame(data)
            
            sentences = []
            for _, row in df.iterrows():
                row_items = []
                for col, value in row.items():
                    if pd.notna(value) and str(value).strip():
                        col_name = str(col).strip()
                        if col_name and col_name != '':
                            row_items.append(f"{col_name}: {value}")
                        else:
                            row_items.append(str(value))
                if row_items:
                    row_text = "; ".join(row_items) + "."
                    sentences.append(row_text)
            return "\n".join(sentences)
        except Exception as e:
            return f"Table content: {str(table)}"

    def ingest_pdf(self, pdf_path: str):
        try:
            print(f"Ingesting PDF from {pdf_path}")
            elements = partition_pdf(
                pdf_path,
                detect_tables=True,
                infer_table_structure=True,
                strategy="hi_res"
            )
            print(f"Found {len(elements)} elements in the PDF")
            self._process_elements(elements, pdf_path)
        except Exception as e:
            print(f"Error ingesting PDF: {str(e)}")
            import traceback
            traceback.print_exc()

    def _convert_ppt_to_pdf(self, ppt_path: str) -> str:
        try:
            file_path = Path(ppt_path)
            pdf_filename = file_path.stem + ".pdf"
            output_dir = os.path.join(os.path.dirname(os.path.dirname(ppt_path)), "pdf")
            os.makedirs(output_dir, exist_ok=True)
            pdf_path = os.path.join(output_dir, pdf_filename)
            system_os = platform.system()
            if system_os == "Windows":
                subprocess.run([
                    'soffice', '--headless', '--convert-to', 'pdf', 
                    '--outdir', output_dir, str(file_path)
                ], check=True)
            elif system_os in ["Linux", "Darwin"]:
                libreoffice_cmd = "libreoffice" if system_os == "Linux" else "soffice"
                subprocess.run([
                    libreoffice_cmd, '--headless', '--convert-to', 'pdf',
                    '--outdir', output_dir, str(file_path)
                ], check=True)
            else:
                raise Exception(f"Unsupported OS: {system_os}")
            if os.path.exists(pdf_path):
                print(f"PDF file exists at: {pdf_path}")
                return pdf_path
            current_dir_pdf = os.path.join(os.path.dirname(ppt_path), pdf_filename)
            if os.path.exists(current_dir_pdf):
                print(f"PDF created in current directory, moving to output dir")
                shutil.move(current_dir_pdf, pdf_path)
                return pdf_path
            print(f"PDF file not found after conversion")
            return None
        except Exception as e:
            print(f"Error converting PPT to PDF: {str(e)}")
            return self._create_pdf_from_ppt_text(ppt_path)

    def _create_pdf_from_ppt_text(self, ppt_path: str) -> str:
        try:
            pdf_filename = os.path.splitext(os.path.basename(ppt_path))[0] + ".pdf"
            output_dir = os.path.join(os.path.dirname(os.path.dirname(ppt_path)), "pdf")
            os.makedirs(output_dir, exist_ok=True)
            pdf_path = os.path.join(output_dir, pdf_filename)
            is_pptx = ppt_path.lower().endswith('.pptx')
            slides_content = []
            if is_pptx:
                presentation = Presentation(ppt_path)
                for i, slide in enumerate(presentation.slides):
                    slide_content = {
                        'number': i+1,
                        'title': slide.shapes.title.text if slide.shapes.title else f"Slide {i+1}",
                        'content': []
                    }
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            slide_content['content'].append(shape.text)
                    slides_content.append(slide_content)
            else:
                slides_content = self._extract_text_from_binary(ppt_path)
            pdf_writer = PdfWriter()
            for slide in slides_content:
                packet = io.BytesIO()
                c = canvas.Canvas(packet, pagesize=letter)
                c.setFont("Helvetica-Bold", 16)
                title = slide.get('title', f"Slide {slide.get('number', 0)}")
                c.drawString(72, 750, title)
                c.setFont("Helvetica", 12)
                y_position = 720
                for text_block in slide.get('content', []):
                    lines = text_block.split('\n')
                    for line in lines:
                        if len(line) > 80:
                            parts = [line[i:i+80] for i in range(0, len(line), 80)]
                            for part in parts:
                                c.drawString(72, y_position, part)
                                y_position -= 15
                                if y_position < 72:
                                    c.showPage()
                                    y_position = 750
                        else:
                            c.drawString(72, y_position, line)
                            y_position -= 15
                            if y_position < 72:
                                c.showPage()
                                y_position = 750
                    y_position -= 10
                c.save()
                packet.seek(0)
                new_pdf = PdfReader(packet)
                for page_num in range(len(new_pdf.pages)):
                    pdf_writer.add_page(new_pdf.pages[page_num])
            if len(slides_content) == 0:
                packet = io.BytesIO()
                c = canvas.Canvas(packet, pagesize=letter)
                filename = os.path.basename(ppt_path)
                c.setFont("Helvetica-Bold", 16)
                c.drawString(72, 750, f"PowerPoint: {filename}")
                c.setFont("Helvetica", 12)
                c.drawString(72, 720, "No slide content could be extracted from this presentation.")
                c.drawString(72, 700, "This is a placeholder page.")
                c.save()
                packet.seek(0)
                new_pdf = PdfReader(packet)
                pdf_writer.add_page(new_pdf.pages[0])
            with open(pdf_path, 'wb') as output_file:
                pdf_writer.write(output_file)
            print(f"Successfully created PDF from extracted text at {pdf_path}")
            return pdf_path
        except Exception as e:
            print(f"Error creating PDF from PPT text: {str(e)}")
            return None

    def _extract_text_from_binary(self, ppt_path: str) -> List[Dict]:
        try:
            print(f"Extracting text directly from PowerPoint binary: {ppt_path}")
            filename = os.path.basename(ppt_path)
            with open(ppt_path, 'rb') as f:
                content = f.read()
            text_chunks = re.findall(b'[a-zA-Z0-9 .,;:\'\"\n\r\t-]{4,}', content)
            extracted_texts = []
            for chunk in text_chunks:
                try:
                    decoded = chunk.decode('utf-8', errors='ignore')
                    if len(decoded) > 10 and not all(c.isdigit() for c in decoded):
                        extracted_texts.append(decoded)
                except:
                    pass
            slides = []
            current_slide_content = []
            current_slide_number = 1
            for text in extracted_texts:
                is_title = (
                    text.startswith("Slide") or 
                    "Title:" in text or 
                    len(text) < 50 and text.strip() and not text.endswith(":")
                )
                if is_title and current_slide_content:
                    slides.append({
                        'number': current_slide_number,
                        'title': f"Slide {current_slide_number}",
                        'content': current_slide_content
                    })
                    current_slide_number += 1
                    current_slide_content = [text]
                else:
                    current_slide_content.append(text)
            if current_slide_content:
                slides.append({
                    'number': current_slide_number,
                    'title': f"Slide {current_slide_number}",
                    'content': current_slide_content
                })
            print(f"Organized binary mesmas em {len(slides)} slides")
            return slides
        except Exception as e:
            print(f"Error extracting text from binary: {str(e)}")
            return []

    def ingest_ppt(self, ppt_path: str) -> bool:
        try:
            print(f"Ingesting PowerPoint from {ppt_path}")
            pdf_path = self._convert_ppt_to_pdf(ppt_path)
            if pdf_path and os.path.exists(pdf_path):
                print(f"Successfully converted to PDF: {pdf_path}")
                self.ingest_pdf(pdf_path)
                success = len(self.vector_store.documents) > 0
                if success:
                    print(f"Successfully processed PPT as PDF with {len(self.vector_store.documents)} sections")
                    return True
                else:
                    print("PPT to PDF conversion succeeded but no content was extracted")
            print("Falling back to direct PowerPoint processing")
            is_pptx = ppt_path.lower().endswith('.pptx')
            if is_pptx:
                presentation = Presentation(ppt_path)
                elements = []
                for i, slide in enumerate(presentation.slides):
                    slide_text = ""
                    if slide.shapes.title and slide.shapes.title.text:
                        slide_text += f"Title: {slide.shapes.title.text}\n\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            slide_text += f"{shape.text}\n"
                    if slide_text.strip():
                        elements.append({
                            'type': 'Slide',
                            'content': slide_text.strip(),
                            'slide_number': i + 1
                        })
                print(f"Extracted {len(elements)} slides with text content")
                doc_count = self._process_slides(elements, ppt_path)
                return doc_count > 0
            else:
                try:
                    import textract
                    text = textract.process(ppt_path).decode('utf-8')
                    if len(text.strip()) > 100:
                        print(f"Extracted {len(text)} characters of text using textract")
                        slide_sections = re.split(r'\n\s*\n+|\r\n\s*\r\n+', text)
                        elements = []
                        for i, section in enumerate(slide_sections):
                            if len(section.strip()) > 10:
                                elements.append({
                                    'type': 'Slide',
                                    'content': section.strip(),
                                    'slide_number': i + 1
                                })
                        print(f"Split content into {len(elements)} sections")
                        doc_count = self._process_slides(elements, ppt_path)
                        return doc_count > 0
                except:
                    pass
                try:
                    result = subprocess.run(['catppt', ppt_path], capture_output=True, text=True)
                    text = result.stdout
                    if len(text.strip()) > 100:
                        print(f"Extracted {len(text)} characters of text using catppt")
                        slide_sections = re.split(r'\n\s*\n+|\r\n\s*\r\n+', text)
                        elements = []
                        for i, section in enumerate(slide_sections):
                            if len(section.strip()) > 10:
                                elements.append({
                                    'type': 'Slide',
                                    'content': section.strip(),
                                    'slide_number': i + 1
                                })
                        print(f"Split content into {len(elements)} sections")
                        doc_count = self._process_slides(elements, ppt_path)
                        return doc_count > 0
                except:
                    pass
                return self._process_ppt_as_text(ppt_path)
        except Exception as e:
            print(f"Error ingesting PowerPoint: {str(e)}")
            return False

    def _process_ppt_as_text(self, ppt_path: str) -> bool:
        try:
            print(f"Extracting text directly from PowerPoint binary: {ppt_path}")
            filename = os.path.basename(ppt_path)
            with open(ppt_path, 'rb') as f:
                content = f.read()
            text_chunks = re.findall(b'[a-zA-Z0-9 .,;:\'\"\n\r\t-]{4,}', content)
            extracted_texts = []
            for chunk in text_chunks:
                try:
                    decoded = chunk.decode('utf-8', errors='ignore')
                    if len(decoded) > 10 and not all(c.isdigit() for c in decoded):
                        extracted_texts.append(decoded)
                except:
                    pass
            sections = []
            current_section = ""
            for text in extracted_texts:
                if len(current_section) > 500 or text.startswith("Title:") or text.startswith("Slide:"):
                    if current_section.strip():
                        sections.append(current_section.strip())
                    current_section = text
                else:
                    if current_section:
                        current_section += "\n" + text
                    else:
                        current_section = text
            if current_section.strip():
                sections.append(current_section.strip())
            elements = []
            for i, section in enumerate(sections):
                if len(section.strip()) > 20:
                    elements.append({
                        'type': 'Slide',
                        'content': section.strip(),
                        'slide_number': i + 1
                    })
            print(f"Extracted {len(elements)} sections from PowerPoint binary")
            if len(elements) < 2:
                print("Adding a minimal document with the file information")
                elements.append({
                    'type': 'Slide',
                    'content': f"PowerPoint presentation: {filename}. Limited text could be extracted.",
                    'slide_number': 1
                })
            doc_count = self._process_slides(elements, ppt_path)
            return doc_count > 0
        except Exception as e:
            print(f"Error in direct PowerPoint text extraction: {str(e)}")
            return False

    def _process_slides(self, elements, source_path):
        slides = {}
        for element in elements:
            slide_number = element.get('slide_number', 0)
            if slide_number not in slides:
                slides[slide_number] = []
            slides[slide_number].append(element)
        print(f"Organized content into {len(slides)} slides")
        document_count = 0
        filename = os.path.basename(source_path)
        for slide_number, slide_elements in slides.items():
            try:
                slide_text = ""
                for element in slide_elements:
                    content = element.get('content', '')
                    if content:
                        slide_text += content + "\n\n"
                slide_text = slide_text.strip()
                if not slide_text:
                    continue
                title = f"Slide {slide_number} - {filename}"
                first_line = slide_text.split('\n')[0] if '\n' in slide_text else ""
                if first_line.startswith("Title:"):
                    title = first_line.replace("Title:", "").strip()
                    slide_text = slide_text[len(first_line):].strip()
                if len(slide_text) > 20000:
                    slide_text = slide_text[:20000] + "... (content truncated)"
                title_embedding = self.get_embedding(title)
                text_embedding = self.get_embedding(slide_text)
                self.vector_store.add_document(
                    title=title,
                    text=slide_text,
                    title_embedding=title_embedding,
                    text_embedding=text_embedding,
                    keywords=["powerpoint", "presentation", "slide", f"slide-{slide_number}"],
                    entities={},
                    metadata={'source': source_path, 'slide': slide_number}
                )
                document_count += 1
            except Exception as slide_error:
                print(f"Error processing slide {slide_number}: {str(slide_error)}")
        print(f"Added {document_count} slides to vector store")
        return document_count

    def _process_elements(self, elements, source_path):
        processed_content = []
        current_section = {'title': '', 'text': []}
        sections_count = 0
        for i, element in enumerate(elements):
            try:
                if isinstance(element, Title):
                    if current_section['title'] or current_section['text']:
                        self._process_section(current_section, source_path)
                        sections_count += 1
                    current_section = {'title': str(element), 'text': []}
                    print(f"Found title: {str(element)[:50]}...")
                elif isinstance(element, Table):
                    table_text = self._convert_table_to_text(element)
                    current_section['text'].append(table_text)
                    print(f"Added table content ({len(table_text)} chars)")
                elif isinstance(element, NarrativeText):
                    current_section['text'].append(str(element))
                    print(f"Added text content ({len(str(element))} chars)")
            except Exception as e:
                print(f"Error processing element {i}: {str(e)}")
        if current_section['title'] or current_section['text']:
            self._process_section(current_section, source_path)
            sections_count += 1
        print(f"Processed {sections_count} sections from the document")
        print(f"Total documents in vector store: {len(self.vector_store.documents)}")

    def _process_section(self, section: Dict, source: str):
        try:
            title = section.get('title', '').strip()
            raw_text = ' '.join([str(t) for t in section.get('text', []) if t])
            text = raw_text.strip()
            if not title and not text:
                print("Skipping completely empty section")
                return False
            if not title:
                text_lines = text.split('\n')
                if text_lines and len(text_lines[0]) > 3:
                    title = text_lines[0][:50]
                else:
                    title = "Untitled Section"
                print(f"Using generated title: {title}")
            if len(text) < 10:
                print(f"Warning: Section has very short content ({len(text)} chars)")
                if title and len(title) > len(text):
                    text = f"{title}. {text}"
                    print(f"Using title as content, new length: {len(text)} chars")
            print(f"Processing section: '{title[:50]}...' with {len(text)} chars of text")
            title_embedding = self.get_embedding(title)
            text_embedding = self.get_embedding(text)
            keywords = self.extract_keywords(text)
            entities = self.extractor.extract_entities(text)
            print(f"Extracted {len(keywords)} keywords and {sum(len(v) for v in entities.values())} entities")
            
            # Print chunked information even if storage fails
            print(f"\n=== Chunked Information (Pre-Storage) ===")
            print(f"Title: {title}")
            print(f"Text (Full Content, {len(text)} chars):\n{text}")
            print(f"Keywords: {', '.join(keywords)}")
            print(f"Entities: {entities}")
            print(f"Metadata: {{'source': '{source}'}}")
            print("=" * 80)

            # Debug: Print arguments
            print("Debug: Arguments for add_document:")
            print(f"  title: {title[:50]}... ({len(title)} chars)")
            print(f"  text: {text[:50]}... ({len(text)} chars)")
            print(f"  title_embedding shape: {title_embedding.shape}")
            print(f"  text_embedding shape: {text_embedding.shape}")
            print(f"  keywords: {keywords}")
            print(f"  entities: {entities}")
            print(f"  metadata: {{'source': '{source}'}}")

            doc_id = self.vector_store.add_document(
                title=title,
                text=text,
                title_embedding=title_embedding,
                text_embedding=text_embedding,
                keywords=keywords,
                entities=entities,
                metadata={'source': source}
            )
            self.chunk_count += 1
            print(f"\nStoring Chunk {self.chunk_count} (ID: {doc_id})")
            print(f"Chunk Text: {text[:500] + '...' if len(text) > 500 else text}")
            print("-" * 50)
            print(f"Added document with ID: {doc_id}")
            return True
        except TypeError as e:
            print(f"TypeError in processing section: {str(e)}")
            import traceback
            traceback.print_exc()
            return False
        except Exception as e:
            print(f"Unexpected error in processing section: {str(e)}")
            import traceback
            traceback.print_exc()
            return False

    def _remove_think_tags(self, text: str) -> str:
        pattern = r'<think>.*?</think>'
        cleaned_text = re.sub(pattern, '', text, flags=re.DOTALL)
        return cleaned_text.strip()

    def _classify_question_type(self, question: str) -> str:
        question_lower = question.lower()
        if re.search(r'(what is|define|meaning of|definition of)', question_lower):
            return "definition"
        elif re.search(r'(list|what are|enumerate)', question_lower):
            return "list"
        elif re.search(r'(summarize|summary|overview|brief|describe)', question_lower):
            return "summary"
        elif re.search(r'(compare|difference between|versus|vs\.)', question_lower):
            return "comparison"
        elif re.search(r'(who|whom|whose)', question_lower):
            return "person"
        elif re.search(r'(when|what time|what date)', question_lower):
            return "temporal"
        elif re.search(r'(where|location|place)', question_lower):
            return "location"
        elif re.search(r'(why|reason|cause)', question_lower):
            return "explanation"
        elif re.search(r'(how many|how much|count|quantity)', question_lower):
            return "quantity"
        else:
            return "general"

    def _extract_question_keywords(self, question: str) -> List[str]:
        question_words = {"what", "when", "where", "who", "whom", "which", "why", "how", 
                         "is", "are", "was", "were", "will", "would", "should", "could", 
                         "can", "do", "does", "did", "has", "have", "had", "the", "a", "an", 
                         "in", "on", "at", "to", "for", "with", "about", "of"}
        doc = nlp(question)
        keywords = []
        for ent in doc.ents:
            keywords.append(ent.text.lower())
        for token in doc:
            if token.pos_ in ["NOUN", "PROPN"] and token.text.lower() not in question_words:
                keywords.append(token.text.lower())
            if token.pos_ == "ADJ" and len(token.text) > 3:
                keywords.append(token.text.lower())
        return list(set(keywords))

    def _generate_non_llm_response(self, question: str, results: List[DocumentEntry]) -> Tuple[str, bool]:
        if not results:
            return "No relevant information found.", False
        question_type = self._classify_question_type(question)
        if question_type == "definition":
            match = re.search(r'(?:what is|define|meaning of|definition of)\s+([^?\.]+)', question.lower())
            if match:
                term = match.group(1).strip()
                for doc in results:
                    definition_patterns = [
                        rf'{term}\s+is\s+defined\s+as\s+([^\.]+)',
                        rf'{term}\s+refers\s+to\s+([^\.]+)',
                        rf'{term}\s+means\s+([^\.]+)',
                        rf'{term}\s+is\s+(?:a|an|the)\s+([^\.]+)'
                    ]
                    for pattern in definition_patterns:
                        matches = re.search(pattern, doc.text.lower())
                        if matches:
                            definition = matches.group(1).strip()
                            return self.templates.generate_response(
                                "definition", 
                                entity=term.capitalize(), 
                                definition=definition
                            ), True
        if question_type == "list":
            match = re.search(r'(?:list|what are|enumerate)\s+([^?\.]+)', question.lower())
            topic = match.group(1).strip() if match else ""
            list_items = []
            for doc in results:
                sentences = sent_tokenize(doc.text)
                for sentence in sentences:
                    if (topic in sentence.lower() or 
                        any(kw in sentence.lower() for kw in self._extract_question_keywords(question))):
                        list_items.append(f"• {sentence}")
                if len(list_items) >= 5:
                    break
            if list_items:
                return self.templates.generate_response(
                    "list",
                    topic=topic,
                    items="\n".join(list_items[:7])
                ), True
        if question_type == "summary":
            match = re.search(r'(?:summarize|summary of|overview of|brief on|describe)\s+([^?\.]+)', question.lower())
            topic = match.group(1).strip() if match else ""
            summary = self.extractive_generator.generate_answer(question, results, max_sentences=5)
            if summary and summary != "No relevant information found.":
                return self.templates.generate_response(
                    "summary",
                    topic=topic if topic else "the topic",
                    summary=summary
                ), True
        if question_type in ["person", "temporal", "location"]:
            entity_types = {
                "person": ["PERSON", "ORG"],
                "temporal": ["DATE", "TIME"],
                "location": ["GPE", "LOC"]
            }
            relevant_entities = []
            for doc in results:
                if doc.entities:
                    for ent_type in entity_types.get(question_type, []):
                        if ent_type in doc.entities:
                            relevant_entities.extend(doc.entities[ent_type])
            if relevant_entities:
                if question_type == "person":
                    return f"The relevant people/organizations are: {', '.join(set(relevant_entities)[:5])}", True
                elif question_type == "temporal":
                    return f"The relevant dates/times are: {', '.join(set(relevant_entities)[:5])}", True
                elif question_type == "location":
                    return f"The relevant locations are: {', '.join(set(relevant_entities)[:5])}", True
        answer = self.extractive_generator.generate_answer(question, results)
        if answer and answer != "No relevant information found.":
            return self.templates.generate_response(
                "fallback",
                content=answer
            ), True
        return "", False

    def _generate_llm_response(self, prompt, max_retries=3):
        if not self.use_llm or not self.generation_model:
            print("LLM not available for generation")
            return None
        for attempt in range(max_retries):
            try:
                print(f"Generating LLM response (attempt {attempt + 1}/{max_retries})")
                response = self.generation_model(prompt)
                if not response or len(response.strip()) < 5:
                    print(f"Empty or very short response received: '{response}'")
                    if attempt < max_retries - 1:
                        print(f"Retrying LLM generation...")
                        time.sleep(2 ** attempt)
                        continue
                    else:
                        print("Maximum retries reached with empty responses")
                        return None
                return response
            except Exception as e:
                print(f"Error in LLM generation attempt {attempt + 1}: {str(e)}")
                if attempt < max_retries - 1:
                    retry_delay = 2 ** attempt
                    print(f"Waiting {retry_delay}s before retry...")
                    time.sleep(retry_delay)
                else:
                    print(f"LLM generation failed after {max_retries} attempts")
                    return None
        return None

    def query(self, question: str, k: int = 5, force_llm: bool = False) -> str:
        try:
            if self._is_developer_question(question):
                return self._get_developer_info()
            question_keywords = self._extract_question_keywords(question)
            keyword_results = self.vector_store.keyword_search(question_keywords, k=k)
            if len(keyword_results) < k:
                try:
                    tfidf_results = self.vector_store.tfidf_search(question, k=k)
                    seen_ids = set(doc.id for doc in keyword_results)
                    for doc in tfidf_results:
                        if doc.id not in seen_ids:
                            keyword_results.append(doc)
                            seen_ids.add(doc.id)
                            if len(keyword_results) >= k:
                                break
                except Exception as e:
                    print(f"Error in TF-IDF search: {str(e)}")
            if len(keyword_results) < k:
                try:
                    query_embedding = self.get_embedding(question)
                    embedding_results = self.vector_store.search(query_embedding, k=k)
                    seen_ids = set(doc.id for doc in keyword_results)
                    for doc in embedding_results:
                        if doc.id not in seen_ids:
                            keyword_results.append(doc)
                            seen_ids.add(doc.id)
                            if len(keyword_results) >= k:
                                break
                except Exception as e:
                    print(f"Error in embedding search: {str(e)}")
            results = keyword_results[:min(k, len(keyword_results))]
            similarities = []
            if results:
                try:
                    query_embedding = self.get_embedding(question)
                    for doc in results:
                        try:
                            if hasattr(doc, 'text_embedding') and doc.text_embedding is not None:
                                sim = self._calculate_cosine_similarity(query_embedding, doc.text_embedding)
                                similarities.append((doc, sim))
                            elif hasattr(doc, 'embedding') and doc.embedding is not None:
                                sim = self._calculate_cosine_similarity(query_embedding, doc.embedding)
                                similarities.append((doc, sim))
                            else:
                                similarities.append((doc, -1))
                        except Exception as e:
                            print(f"Error calculating similarity for document: {str(e)}")
                            similarities.append((doc, -1))
                    similarities.sort(key=lambda x: x[1], reverse=True)
                    results = [doc for doc, _ in similarities]
                except Exception as e:
                    print(f"Error calculating similarities: {str(e)}")
            if not results:
                return self._format_response("I couldn't find any relevant information in the document to answer your question. Please try asking something else about the content of the document.")
            similarity_info = self._format_similarity_info(similarities)
            if self.method == "rule" and not force_llm:
                non_llm_response, success = self._generate_non_llm_response(question, results)
                if success:
                    return self._format_response(non_llm_response + similarity_info)
                return self._format_response("No relevant rule-based answer found." + similarity_info)
            elif self.method == "extractive" and not force_llm:
                try:
                    answer = self.extractive_generator.generate_answer(question, results)
                    return self._format_response(answer + similarity_info)
                except Exception as e:
                    print(f"Error in extractive generation: {str(e)}")
                    return self._format_response("I couldn't generate an answer using the extractive method. Please try a different approach." + similarity_info)
            elif self.method == "extraction" and not force_llm:
                all_entities = {}
                for doc in results:
                    if hasattr(doc, 'entities') and doc.entities:
                        for entity_type, entities in doc.entities.items():
                            if entity_type not in all_entities:
                                all_entities[entity_type] = set()
                            all_entities[entity_type].update(entities)
                if all_entities:
                    response_parts = ["Key information extracted from the documents:"]
                    for entity_type, entities in all_entities.items():
                        if len(entities) > 0:
                            entity_list = ", ".join(list(entities)[:5])
                            response_parts.append(f"{entity_type}: {entity_list}")
                    return self._format_response("\n".join(response_parts) + similarity_info)
                return self._format_response("No relevant entities found in the documents." + similarity_info)
            elif self.method == "tfidf" and not force_llm:
                try:
                    tfidf_results = self.vector_store.tfidf_search(question, k=k)
                    if tfidf_results and len(tfidf_results) > 0:
                        top_result = tfidf_results[0]
                        response = f"Top match: {top_result.title}\n\n{top_result.text[:500]}..."
                        return self._format_response(response + similarity_info)
                    return self._format_response("No relevant documents found using TF-IDF search." + similarity_info)
                except Exception as e:
                    print(f"Error in TF-IDF search processing: {str(e)}")
                    return self._format_response("Error processing TF-IDF search results." + similarity_info)
            elif self.use_llm:
                try:
                    context_parts = []
                    for i, (doc, sim) in enumerate(similarities[:min(k, len(similarities))]):
                        if sim < 0 or sim >= 0.1:
                            relevance_str = f"{sim:.2f}" if sim >= 0 else "N/A"
                            context_parts.append(f"Section {i+1} (Relevance: {relevance_str}):\nTitle: {doc.title}\nContent: {doc.text[:1000]}...")
                    context = "\n\n".join(context_parts)
                    prompt = f"""Answer ONLY based on the information in these sections. If you can't find the answer in the text below, say "I don't have information about that in the document."

{context}

Question: {question}

Answer:"""
                    response = self._generate_llm_response(prompt)
                    if response is None:
                        print("LLM generation failed, falling back to extractive answer")
                        extracted_answer = self.extractive_generator.generate_answer(question, results)
                        return self._format_response(extracted_answer + similarity_info)
                    cleaned_response = self._remove_think_tags(response)
                    if len(cleaned_response.strip()) < 10:
                        print("LLM response was too short, falling back to extractive answer")
                        extracted_answer = self.extractive_generator.generate_answer(question, results)
                        return self._format_response(extracted_answer + similarity_info)
                    return self._format_response(cleaned_response + similarity_info)
                except Exception as e:
                    print(f"Error generating LLM response: {str(e)}")
                    print("Falling back to extractive answer generation")
                    extracted_answer = self.extractive_generator.generate_answer(question, results)
                    return self._format_response(extracted_answer + similarity_info)
            else:
                try:
                    answer = self.extractive_generator.generate_answer(question, results)
                    return self._format_response(answer + similarity_info)
                except Exception as e:
                    print(f"Error in default extractive generation: {str(e)}")
                    return self._format_response("I couldn't generate an answer using the extractive method. Please try a different approach." + similarity_info)
        except Exception as e:
            print(f"Unexpected error in query processing: {str(e)}")
            return self._format_response("I encountered an unexpected error while processing your question. Please try again with a different question.")

    def _is_developer_question(self, question: str) -> bool:
        question_lower = question.lower()
        developer_keywords = ["developer", "creator", "author", "made", "created", "built", "developed", 
                              "who are you", "who created", "who made", "who built", "who developed", 
                              "your maker", "who designed"]
        return any(keyword in question_lower for keyword in developer_keywords)
        
    def _get_developer_info(self) -> str:
        info = self.developer_info
        response = f"I was developed by {info['name']}, {info['title']} at {info['company']}.\n\n"
        response += f"You can find more about the developer at:\n"
        response += f"- GitHub: {info['github']}\n"
        response += f"- ResearchGate: {info['researchgate']}\n"
        response += f"- Email: {info['email']}"
        return self._format_response(response)
        
    def _format_similarity_info(self, similarities) -> str:
        if not similarities:
            return ""
        info = "\n\nRelevance scores:\n"
        for i, (doc, sim) in enumerate(similarities[:3]):
            if i >= 3:
                break
            title = doc.title[:30] + "..." if len(doc.title) > 30 else doc.title
            score = f"{sim:.2f}" if sim >= 0 else "N/A"
            info += f"- {title}: {score}\n"
        return info
        
    def _format_response(self, response: str) -> str:
        model_type = "LLM" if self.use_llm else "Non-LLM"
        model_info = f"Response generated from {self.model_name} ({model_type})\n\n"
        return model_info + response

    def _calculate_cosine_similarity(self, embedding1, embedding2):
        try:
            if not isinstance(embedding1, np.ndarray):
                embedding1 = np.array(embedding1)
            if not isinstance(embedding2, np.ndarray):
                embedding2 = np.array(embedding2)
            norm1 = np.linalg.norm(embedding1)
            norm2 = np.linalg.norm(embedding2)
            if norm1 == 0 or norm2 == 0:
                return 0.0
            return float(np.dot(embedding1, embedding2) / (norm1 * norm2))
        except Exception as e:
            print(f"Error calculating similarity: {str(e)}")
            return 0.0

    def save(self, path_prefix: str):
        self.vector_store.save(path_prefix)

    @classmethod
    def load(cls, path_prefix: str, model_name: str = "deepseek-r1:14b", use_llm: bool = True, method: str = "hybrid"):
        rag = cls(model_name, use_llm=use_llm, method=method)
        rag.vector_store = ParallelVectorStore.load(path_prefix)
        return rag

    def ingest_doc(self, doc_path: str):
        try:
            print(f"Ingesting Word document from {doc_path}")
            if not os.path.exists(doc_path):
                print(f"Error: Word document not found at {doc_path}")
                return
            file_extension = os.path.splitext(doc_path)[1].lower()
            if file_extension == '.docx':
                elements = partition_docx(
                    doc_path,
                    detect_tables=True,
                    infer_table_structure=True
                )
            else:
                elements = partition_doc(
                    doc_path,
                    detect_tables=True,
                    infer_table_structure=True
                )
            if not elements:
                print("Warning: No elements extracted from Word document")
                filename = os.path.basename(doc_path)
                elements = [Title(text=f"Document: {filename}")]
            print(f"Found {len(elements)} elements in the Word document")
            self._process_elements(elements, doc_path)
        except Exception as e:
            print(f"Error ingesting Word document: {str(e)}")
            import traceback
            traceback.print_exc()

    def ingest_spreadsheet(self, file_path: str) -> bool:
        try:
            print(f"Ingesting spreadsheet from {file_path}")
            pdf_path = self._convert_spreadsheet_to_pdf(file_path)
            if pdf_path and os.path.exists(pdf_path):
                print(f"Successfully converted to PDF: {pdf_path}")
                self.ingest_pdf(pdf_path)
                success = len(self.vector_store.documents) > 0
                if success:
                    print(f"Successfully processed spreadsheet as PDF with {len(self.vector_store.documents)} sections")
                    return True
                else:
                    print("Spreadsheet to PDF conversion succeeded but no content was extracted")
            return False
        except Exception as e:
            print(f"Error ingesting spreadsheet: {str(e)}")
            return False

    def _convert_spreadsheet_to_pdf(self, file_path: str) -> str:
        try:
            import pandas as pd
            from reportlab.lib import colors
            from reportlab.lib.pagesizes import letter, landscape
            from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
            from reportlab.lib.units import inch
            print(f"Starting spreadsheet to PDF conversion for {file_path}")
            pdf_filename = os.path.splitext(os.path.basename(file_path))[0] + ".pdf"
            output_dir = os.path.join(os.path.dirname(os.path.dirname(file_path)), "pdf")
            os.makedirs(output_dir, exist_ok=True)
            pdf_path = os.path.join(output_dir, pdf_filename)
            file_ext = os.path.splitext(file_path)[-1].lower()
            if file_ext == '.csv':
                print(f"Reading CSV file: {file_path}")
                df = pd.read_csv(file_path)
            elif file_ext in ['.xls', '.xlsx']:
                print(f"Reading Excel file: {file_path}")
                df = pd.read_excel(file_path, engine='openpyxl')
            else:
                raise ValueError(f"Unsupported file format: {file_ext}. Please provide a CSV or Excel file.")
            print(f"Successfully loaded dataframe with {len(df)} rows and {len(df.columns)} columns")
            max_rows = min(1000, len(df))
            if max_rows < len(df):
                print(f"Limiting output to {max_rows} rows out of {len(df)}")
                df = df.iloc[:max_rows]
            data = [df.columns.tolist()]
            for i, row in df.iterrows():
                row_data = []
                for item in row:
                    if pd.isna(item):
                        row_data.append('')
                    else:
                        item_str = str(item)
                        if len(item_str) > 100:
                            item_str = item_str[:97] + '...'
                        row_data.append(item_str)
                data.append(row_data)
            doc = SimpleDocTemplate(
                pdf_path,
                pagesize=landscape(letter),
                rightMargin=0.5*inch,
                leftMargin=0.5*inch,
                topMargin=0.5*inch,
                bottomMargin=0.5*inch
            )
            num_cols = len(data[0])
            available_width = landscape(letter)[0] - inch
            col_width = min(1.5*inch, available_width / num_cols)
            col_widths = [col_width] * num_cols
            table = Table(data, colWidths=col_widths)
            style = TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
            ])
            table.setStyle(style)
            elements = [table]
            print(f"Building PDF document...")
            doc.build(elements)
            if os.path.exists(pdf_path):
                print(f"PDF successfully created at: {pdf_path}")
                return pdf_path
            else:
                print(f"Error: PDF file was not created at {pdf_path}")
                return None
        except Exception as e:
            print(f"Error converting spreadsheet to PDF: {str(e)}")
            return None

# Usage Example
if __name__ == "__main__":
    # Example paths for different document types
    PDF_PATH = "/home/shahanahmed/Downloads/DocumentForChatbot/AboutBangladesh.pdf"
    PPTX_PATH = "/home/shahanahmed/Documents/presentation.pptx"
    DOCX_PATH = "/home/shahanahmed/Documents/document.docx"

    # Create a RAG system
    rag = RAGSystem(model_name="gemma3:latest", use_llm=True, method="hybrid")
    
    # Choose which file type to process
    document_type = "pdf"  # Change to "pptx" or "docx" to test other document types
    
    if document_type == "pdf":
        rag.ingest_pdf(PDF_PATH)
    elif document_type == "pptx":
        rag.ingest_ppt(PPTX_PATH) 
    elif document_type == "docx":
        rag.ingest_doc(DOCX_PATH)
    
    # Test with a sample question
    question = "Tell me about the document"
    rag.vector_store.print_all_documents()
    # # Hybrid approach
    # print("Hybrid approach (non-LLM with LLM fallback):")
    # rag.method = "hybrid"
    # print(rag.query(question))
    
    # # Rule-based approach
    # print("\nPure non-LLM approach (rule):")
    # rag.method = "rule"
    # print(rag.query(question, force_llm=False))

    # # Extractive approach
    # print("\nPure non-LLM approach (extractive):")
    # rag.method = "extractive"
    # print(rag.query(question, force_llm=False))

    # # Entity extraction approach
    # print("\nPure non-LLM approach (extraction):")
    # rag.method = "extraction"
    # print(rag.query(question, force_llm=False))

    # # TF-IDF approach
    # print("\nPure non-LLM approach (tfidf):")
    # rag.method = "tfidf"
    # print(rag.query(question, force_llm=False))